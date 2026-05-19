# -*- coding: utf-8 -*-
"""Generate a PPTX report on the PEFA paper (WSDM '24)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT_PATH = r"C:\Users\yiwen.liu\Desktop\114-2 IR\report_2\PEFA_report.pptx"
SRC_DIR  = r"C:\Users\yiwen.liu\Desktop\114-2 IR\report_2\src"
FIG1     = os.path.join(SRC_DIR, "figure1.png")  # PEFA-XL & PEFA-XS architecture
FIG3     = os.path.join(SRC_DIR, "figure3.png")  # supervised-data ablation

# ---------- Theme colors ----------
PRIMARY   = RGBColor(0x1F, 0x4E, 0x79)   # deep blue
ACCENT    = RGBColor(0xC0, 0x50, 0x4D)   # warm red
SOFT_BG   = RGBColor(0xF2, 0xF2, 0xF2)
DARK_TXT  = RGBColor(0x22, 0x22, 0x22)
LIGHT_TXT = RGBColor(0xFF, 0xFF, 0xFF)
SUBTLE    = RGBColor(0x59, 0x59, 0x59)
GREEN     = RGBColor(0x4E, 0x7B, 0x3A)
ORANGE    = RGBColor(0xD9, 0x82, 0x2B)

# ---------- Presentation setup ----------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=DARK_TXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, bullets, *, size=16, color=DARK_TXT,
                bold_first_token=False, line_spacing=1.15, bullet_char="•"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # bullet character
        r0 = p.add_run()
        r0.text = f"{bullet_char}  "
        r0.font.name = "Calibri"
        r0.font.size = Pt(size)
        r0.font.bold = True
        r0.font.color.rgb = PRIMARY
        # main content; allow "**bold** rest"
        if bold_first_token and " — " in item:
            head, tail = item.split(" — ", 1)
            rh = p.add_run()
            rh.text = head + " — "
            rh.font.name = "Calibri"
            rh.font.size = Pt(size)
            rh.font.bold = True
            rh.font.color.rgb = color
            rt = p.add_run()
            rt.text = tail
            rt.font.name = "Calibri"
            rt.font.size = Pt(size)
            rt.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.name = "Calibri"
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


def add_header(slide, title, subtitle=None, page=None, total=None):
    add_rect(slide, 0, 0, SW, Inches(0.85), PRIMARY)
    add_rect(slide, 0, Inches(0.85), SW, Inches(0.06), ACCENT)
    add_text(slide, Inches(0.45), Inches(0.12), Inches(11), Inches(0.6),
             title, size=26, bold=True, color=LIGHT_TXT,
             anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.45), Inches(0.95), Inches(11), Inches(0.35),
                 subtitle, size=13, color=SUBTLE, anchor=MSO_ANCHOR.TOP)
    if page is not None and total is not None:
        add_text(slide, Inches(11.8), Inches(0.2),
                 Inches(1.3), Inches(0.5),
                 f"{page} / {total}", size=12, bold=True,
                 color=LIGHT_TXT, align=PP_ALIGN.RIGHT,
                 anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, text="PEFA · WSDM '24 · Paper Reading Report"):
    add_text(slide, Inches(0.3), Inches(7.12), Inches(13), Inches(0.3),
             text, size=10, color=SUBTLE)


def add_table(slide, x, y, w, h, data, *, header_fill=PRIMARY,
              header_color=LIGHT_TXT, body_color=DARK_TXT,
              first_col_bold=True, header_size=12, body_size=11,
              col_widths=None):
    rows, cols = len(data), len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    tbl = tbl_shape.table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(w * cw / total)
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            tf = cell.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = ""
            run = tf.paragraphs[0].add_run()
            run.text = str(data[r][c])
            run.font.name = "Calibri"
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                run.font.bold = True
                run.font.size = Pt(header_size)
                run.font.color.rgb = header_color
            else:
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = SOFT_BG
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.size = Pt(body_size)
                run.font.color.rgb = body_color
                if c == 0 and first_col_bold:
                    run.font.bold = True
                    run.font.color.rgb = PRIMARY
            tf.paragraphs[0].alignment = (PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)
    return tbl


def add_callout(slide, x, y, w, h, title, body, *, color=ACCENT):
    add_rect(slide, x, y, Inches(0.12), h, color)
    add_rect(slide, x + Inches(0.12), y, w - Inches(0.12), h, SOFT_BG)
    add_text(slide, x + Inches(0.25), y + Inches(0.08),
             w - Inches(0.35), Inches(0.4),
             title, size=14, bold=True, color=color)
    add_text(slide, x + Inches(0.25), y + Inches(0.5),
             w - Inches(0.35), h - Inches(0.55),
             body, size=12, color=DARK_TXT)


TOTAL = 15  # 14 base + 1 new "Effect of supervised data" slide for figure3

# ============================================================
# Slide 1 — Cover
# ============================================================
s = prs.slides.add_slide(BLANK)
# background gradient-like layers
add_rect(s, 0, 0, SW, SH, RGBColor(0xF8, 0xFA, 0xFC))
add_rect(s, 0, 0, SW, Inches(2.8), PRIMARY)
add_rect(s, 0, Inches(2.8), SW, Inches(0.08), ACCENT)
# Accent shapes
tri = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                         Inches(10.8), Inches(0), Inches(2.6), Inches(2.8))
tri.fill.solid(); tri.fill.fore_color.rgb = ACCENT
tri.line.fill.background()
tri.rotation = 180

add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.5),
         "Paper Reading Report · 114-2 Information Retrieval",
         size=14, color=LIGHT_TXT)
add_text(s, Inches(0.6), Inches(1.05), Inches(12), Inches(1.2),
         "PEFA: Parameter-Free Adapters",
         size=40, bold=True, color=LIGHT_TXT)
add_text(s, Inches(0.6), Inches(1.85), Inches(12), Inches(0.6),
         "for Large-scale Embedding-based Retrieval Models",
         size=22, color=LIGHT_TXT)

# Authors / venue card
add_rect(s, Inches(0.6), Inches(3.4), Inches(12.1), Inches(2.6),
         RGBColor(0xFF, 0xFF, 0xFF), line=PRIMARY)
add_text(s, Inches(0.9), Inches(3.55), Inches(11.5), Inches(0.45),
         "Authors", size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.9), Inches(3.9), Inches(11.5), Inches(0.55),
         "Wei-Cheng Chang, Jyun-Yu Jiang, Jiong Zhang, Mutasem Al-Darabsah, "
         "Choon Hui Teo, Cho-Jui Hsieh, Hsiang-Fu Yu, S. V. N. Vishwanathan",
         size=14, color=DARK_TXT)
add_text(s, Inches(0.9), Inches(4.5), Inches(11.5), Inches(0.4),
         "Affiliation", size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.9), Inches(4.82), Inches(11.5), Inches(0.4),
         "Amazon  ·  UCLA",
         size=14, color=DARK_TXT)
add_text(s, Inches(0.9), Inches(5.25), Inches(5.5), Inches(0.4),
         "Venue", size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.9), Inches(5.55), Inches(5.5), Inches(0.4),
         "WSDM '24 · Mérida, Mexico",
         size=14, color=DARK_TXT)
add_text(s, Inches(7), Inches(5.25), Inches(5.5), Inches(0.4),
         "Reference", size=12, bold=True, color=ACCENT)
add_text(s, Inches(7), Inches(5.55), Inches(5.7), Inches(0.4),
         "arXiv:2312.02429v2  ·  github.com/amzn/pecos",
         size=14, color=DARK_TXT)

add_text(s, Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
         "Presenter: 314581027 劉怡妏    |    Date: 2026-05-19",
         size=12, color=SUBTLE)

# ============================================================
# Slide 2 — Outline
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Outline", "Roadmap of this presentation", page=2, total=TOTAL)

items = [
    ("1", "Motivation", "Why full-parameter fine-tuning is too expensive"),
    ("2", "Preliminaries", "Dense retrieval, ERMs, MIPS & ANN"),
    ("3", "PEFA Framework", "Convex combination of ERM + non-parametric kNN"),
    ("4", "PEFA-XL & PEFA-XS", "Two realizations with different trade-offs"),
    ("5", "Complexity Analysis", "Index size, build time, inference latency"),
    ("6", "Document Retrieval", "Experiments on NQ-320K & Trivia-QA"),
    ("7", "Product Search", "Billion-scale industrial deployment"),
    ("8", "Ablation & Discussion", "Hyper-parameters, related work, takeaways"),
]
cols = 2
rows = 4
cw = Inches(6.05)
rh = Inches(1.25)
x0 = Inches(0.5)
y0 = Inches(1.6)
gap = Inches(0.15)
for i, (num, head, desc) in enumerate(items):
    r, c = i // cols, i % cols
    x = x0 + c * (cw + gap)
    y = y0 + r * (rh + gap)
    add_rect(s, x, y, cw, rh, RGBColor(0xFF, 0xFF, 0xFF), line=PRIMARY)
    add_rect(s, x, y, Inches(1.0), rh, PRIMARY)
    add_text(s, x, y, Inches(1.0), rh,
             num, size=32, bold=True, color=LIGHT_TXT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(1.15), y + Inches(0.15),
             cw - Inches(1.2), Inches(0.5),
             head, size=18, bold=True, color=PRIMARY)
    add_text(s, x + Inches(1.15), y + Inches(0.62),
             cw - Inches(1.2), Inches(0.6),
             desc, size=12, color=DARK_TXT)
add_footer(s)

# ============================================================
# Slide 3 — Motivation
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "1 · Motivation",
           "Fine-tuning ERMs at industry scale is painful", page=3, total=TOTAL)

# Pain points (left)
add_text(s, Inches(0.5), Inches(1.5), Inches(6.3), Inches(0.4),
         "Pain points of full-parameter fine-tuning",
         size=18, bold=True, color=PRIMARY)
pains = [
    "Multi-stage pipeline — pre-train → 1st-stage FT (random/BM25 neg) → 2nd-stage FT (hard neg) → 3rd-stage distillation",
    "Compute cost — billions of (query, product) pairs → thousands of GPU hours",
    "Gradient access required — not feasible for black-box LLMs (GPT-3, etc.)",
    "Tail queries & labels — even SoTA bi-encoders struggle on long-tail",
]
add_bullets(s, Inches(0.5), Inches(1.95), Inches(6.3), Inches(4.5),
            pains, size=13)

# Right: contrasting illustration
add_rect(s, Inches(7.1), Inches(1.5), Inches(5.7), Inches(5.1),
         SOFT_BG)
add_text(s, Inches(7.3), Inches(1.6), Inches(5.4), Inches(0.4),
         "PEFA's promise", size=18, bold=True, color=ACCENT)

# Two big boxes: traditional vs PEFA
add_rect(s, Inches(7.3), Inches(2.1), Inches(5.3), Inches(1.9),
         RGBColor(0xFF, 0xFF, 0xFF), line=SUBTLE)
add_text(s, Inches(7.45), Inches(2.18), Inches(5.0), Inches(0.4),
         "Traditional fine-tuning",
         size=14, bold=True, color=SUBTLE)
add_bullets(s, Inches(7.45), Inches(2.55), Inches(5.0), Inches(1.4),
            ["Backward pass + gradient updates",
             "GPU clusters, days to weeks",
             "Tied to white-box models"],
            size=12, line_spacing=1.1)

add_rect(s, Inches(7.3), Inches(4.15), Inches(5.3), Inches(2.4),
         RGBColor(0xFF, 0xFF, 0xFF), line=ACCENT)
add_text(s, Inches(7.45), Inches(4.23), Inches(5.0), Inches(0.4),
         "PEFA (this paper)",
         size=14, bold=True, color=ACCENT)
add_bullets(s, Inches(7.45), Inches(4.6), Inches(5.0), Inches(1.9),
            ["No backward pass · no gradient",
             "CPU-side ANN index build (hours, not days)",
             "Works on pre-trained AND fine-tuned ERMs",
             "Applies to black-box LLM encoders"],
            size=12, line_spacing=1.1)
add_footer(s)

# ============================================================
# Slide 4 — Preliminaries: Dense retrieval
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "2 · Preliminaries",
           "Dense text retrieval with bi-encoder ERMs",
           page=4, total=TOTAL)

# Equation block
add_rect(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(1.4),
         SOFT_BG)
add_text(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
         "ERM scoring function",
         size=14, bold=True, color=PRIMARY)
add_text(s, Inches(0.7), Inches(1.92), Inches(12), Inches(0.85),
         "f_ERM(q, p; θ) = ⟨ E(q; θ),  E(p; θ) ⟩",
         size=22, bold=True, color=DARK_TXT, font="Cambria Math")
add_text(s, Inches(0.7), Inches(2.5), Inches(12), Inches(0.4),
         "E(·; θ) : X → R^d  ·  inner product / cosine similarity assumed (ℓ2-normalized)",
         size=12, color=SUBTLE)

# Two columns: learning vs inference
y0 = Inches(3.1)
add_rect(s, Inches(0.5), y0, Inches(6.0), Inches(3.7),
         RGBColor(0xFF, 0xFF, 0xFF), line=PRIMARY)
add_text(s, Inches(0.7), y0 + Inches(0.08), Inches(5.6), Inches(0.4),
         "Learning",
         size=16, bold=True, color=PRIMARY)
add_bullets(s, Inches(0.7), y0 + Inches(0.5), Inches(5.6), Inches(3.1),
            ["Training data: relevant pairs D = {(qᵢ, pᵢ)}",
             "Maximize log p_θ(p | q)  (softmax over corpus)",
             "Partition function is intractable → negative sampling: random, BM25, hard mining, distillation",
             "All stages require gradient back-prop"],
            size=12)

add_rect(s, Inches(6.8), y0, Inches(6.0), Inches(3.7),
         RGBColor(0xFF, 0xFF, 0xFF), line=ACCENT)
add_text(s, Inches(7.0), y0 + Inches(0.08), Inches(5.6), Inches(0.4),
         "Inference",
         size=16, bold=True, color=ACCENT)
add_bullets(s, Inches(7.0), y0 + Inches(0.5), Inches(5.6), Inches(3.1),
            ["Maximum Inner Product Search (MIPS) over corpus P",
             "Exact MIPS is O(n) — infeasible for 100M+ corpora",
             "ANN libraries: HNSW, FAISS, ScaNN — sub-linear O(log n)",
             "Index build is cheap (CPU) vs FT (GPU clusters)"],
            size=12)
add_footer(s)

# ============================================================
# Slide 5 — PEFA Framework (core idea)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "3 · PEFA Framework",
           "Convex combination of ERM and non-parametric kNN",
           page=5, total=TOTAL)

# Big equation
add_rect(s, Inches(0.5), Inches(1.45), Inches(12.3), Inches(1.8),
         PRIMARY)
add_text(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
         "Core scoring function (applied only at inference)",
         size=13, bold=True, color=LIGHT_TXT)
add_text(s, Inches(0.7), Inches(1.95), Inches(12), Inches(0.9),
         "f_PEFA(q̂, p_j) = λ · f_ERM(q̂, p_j)  +  (1 − λ) · f_kNN(q̂, p_j)",
         size=24, bold=True, color=LIGHT_TXT, font="Cambria Math",
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.7), Inches(2.78), Inches(12), Inches(0.4),
         "λ ∈ [0,1] :  interpolation coefficient.  ERM parameters θ remain unchanged.",
         size=12, color=LIGHT_TXT, align=PP_ALIGN.CENTER)

# Three property cards
cards = [
    ("Parameter-free",
     "No optimization over ERM weights.\nLearning reduces to building an ANN index of key-value pairs.",
     PRIMARY),
    ("Black-box friendly",
     "Needs only forward embeddings.\nApplies to GPT-3-class LLM encoders with no gradient access.",
     ACCENT),
    ("Orthogonal",
     "Complementary to pre-training, FT, and parameter-efficient methods (Adapter, LoRA, prefix-tuning).",
     GREEN),
]
x0 = Inches(0.5)
y0 = Inches(3.6)
cw = Inches(4.05)
ch = Inches(3.2)
gap = Inches(0.18)
for i, (title, body, col) in enumerate(cards):
    x = x0 + i * (cw + gap)
    add_rect(s, x, y0, cw, ch, RGBColor(0xFF, 0xFF, 0xFF), line=col)
    add_rect(s, x, y0, cw, Inches(0.55), col)
    add_text(s, x, y0, cw, Inches(0.55),
             title, size=16, bold=True, color=LIGHT_TXT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + Inches(0.2), y0 + Inches(0.7),
             cw - Inches(0.4), ch - Inches(0.8),
             body, size=13, color=DARK_TXT)
add_footer(s)

# ============================================================
# Slide 6 — PEFA-XL
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "4a · PEFA-XL  (extra large)",
           "kNN over training queries — query-dependent neighborhood",
           page=6, total=TOTAL)

# Left: definition + equation
add_text(s, Inches(0.5), Inches(1.4), Inches(7.5), Inches(0.4),
         "Idea",
         size=16, bold=True, color=PRIMARY)
add_bullets(s, Inches(0.5), Inches(1.85), Inches(7.5), Inches(2.0),
            ["Find top-k most similar training queries qᵢ ∈ Q to test query q̂",
             "Aggregate the relevant passages of those neighbors (Y_{i,j})",
             "Diagonal gate D_{i,i} = 1 if i ∈ NN(q̂, Q; k); else 0"],
            size=13)

# Equation card
add_rect(s, Inches(0.5), Inches(3.9), Inches(7.5), Inches(1.55),
         SOFT_BG)
add_text(s, Inches(0.7), Inches(3.98), Inches(7.2), Inches(0.4),
         "PEFA-XL scoring (Eq. 6)",
         size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.7), Inches(4.32), Inches(7.2), Inches(1.0),
         "f_XL(q̂, p_j) = λ⟨q̂, p_j⟩ + (1−λ) · Σ_{i ∈ NN(q̂,Q;k)} ⟨q̂, qᵢ⟩ · Y_{i,j}",
         size=15, bold=True, color=DARK_TXT, font="Cambria Math")

add_text(s, Inches(0.5), Inches(5.6), Inches(7.5), Inches(0.4),
         "Normalization",
         size=14, bold=True, color=PRIMARY)
add_text(s, Inches(0.5), Inches(5.95), Inches(7.5), Inches(0.9),
         "D_{i,i} = 1/k so that f_kNN stays bounded by 1, calibrated against f_ERM (cosine in [−1, 1]).",
         size=12, color=DARK_TXT)

# Right: paper Figure 1 (PEFA-XL), cropped from the combined figure1.png
dx = Inches(8.3)
add_rect(s, dx, Inches(1.4), Inches(4.8), Inches(5.4),
         RGBColor(0xFF, 0xFF, 0xFF), line=ACCENT)
add_text(s, dx, Inches(1.45), Inches(4.8), Inches(0.35),
         "Figure 1 of the paper · PEFA-XL",
         size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
pic = s.shapes.add_picture(FIG1, dx + Inches(0.1), Inches(1.85),
                           width=Inches(4.6))
# crop right half away (figure1.png contains both XL and XS side by side)
pic.crop_right = 0.5
# scale down to fit available height
pic.height = Inches(4.6)
pic.width = Inches(4.6)
add_text(s, dx + Inches(0.1), Inches(6.5), Inches(4.6), Inches(0.4),
         "Two ANN indices · TWO ANN searches at inference.",
         size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_footer(s)

# ============================================================
# Slide 7 — PEFA-XS
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "4b · PEFA-XS  (extra small)",
           "Query-independent kNN — fold into a single ANN index",
           page=7, total=TOTAL)

add_text(s, Inches(0.5), Inches(1.4), Inches(7.5), Inches(0.4),
         "Idea",
         size=16, bold=True, color=PRIMARY)
add_bullets(s, Inches(0.5), Inches(1.85), Inches(7.5), Inches(2.0),
            ["Replace NN(q̂, Q; k) with I(p_j, Y) — queries relevant to passage p_j",
             "Neighborhood is now independent of test query q̂",
             "→ interpolation can be pre-computed OFFLINE in embedding space"],
            size=13)

add_rect(s, Inches(0.5), Inches(3.9), Inches(7.5), Inches(1.7),
         SOFT_BG)
add_text(s, Inches(0.7), Inches(3.98), Inches(7.2), Inches(0.4),
         "PEFA-XS scoring (Eq. 8)",
         size=12, bold=True, color=ACCENT)
add_text(s, Inches(0.7), Inches(4.35), Inches(7.2), Inches(1.1),
         "f_XS(q̂, p_j) = ⟨ q̂,  λ · p_j + (1−λ) · Π(Q^T Y_{:,j}) ⟩",
         size=15, bold=True, color=DARK_TXT, font="Cambria Math")
add_text(s, Inches(0.7), Inches(5.15), Inches(7.2), Inches(0.5),
         "Π(·) is ℓ2-normalization onto the unit sphere.",
         size=11, color=SUBTLE)

add_callout(s, Inches(0.5), Inches(5.8), Inches(7.5), Inches(1.2),
            "Why it matters",
            "The interpolated passage embedding is built ONCE at offline indexing — "
            "inference cost is identical to the original ERM. Zero deployment overhead.")

# Right: paper Figure 2 (PEFA-XS), cropped from the combined figure1.png
dx = Inches(8.3)
add_rect(s, dx, Inches(1.4), Inches(4.8), Inches(5.4),
         RGBColor(0xFF, 0xFF, 0xFF), line=PRIMARY)
add_text(s, dx, Inches(1.45), Inches(4.8), Inches(0.35),
         "Figure 2 of the paper · PEFA-XS",
         size=12, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
pic = s.shapes.add_picture(FIG1, dx + Inches(0.1), Inches(1.85),
                           width=Inches(4.6))
# crop left half away to keep only the PEFA-XS panel
pic.crop_left = 0.5
pic.height = Inches(4.6)
pic.width = Inches(4.6)
add_text(s, dx + Inches(0.1), Inches(6.5), Inches(4.6), Inches(0.4),
         "Single ANN index · ONE ANN search · zero extra latency.",
         size=11, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
add_footer(s)

# ============================================================
# Slide 8 — Complexity comparison
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "5 · Complexity Analysis",
           "Time & space at inference, using HNSW as the ANN backbone",
           page=8, total=TOTAL)

data = [
    ["Method",     "Build time",                "Index size",                                    "Inference time"],
    ["ERM",        "O(n log n)",                "O(nd + |E_P|)",                                 "O(log n)"],
    ["PEFA-XS",    "O(n log n)",                "O(nd + |E_P|)",                                 "O(log n)"],
    ["PEFA-XL",    "O(n log n + m log m)",      "O(nd + |E_P| + md + |E_Q| + nnz(Y))",           "O(log n + log m)"],
]
add_table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.4),
          data, body_size=13, header_size=13,
          col_widths=[1.6, 2.6, 5.5, 2.6])

# Key takeaways
add_text(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.4),
         "Key takeaways",
         size=16, bold=True, color=PRIMARY)
add_bullets(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.4),
            ["PEFA-XS has the SAME asymptotic cost as the ERM baseline — drop-in.",
             "PEFA-XL adds a second HNSW over the training-query space Q (size m), often comparable to or larger than n in e-commerce.",
             "Storing Y costs only O(nnz(Y)) — sparse, since Y ∈ {0, 1}^{m×n}.",
             "Index building runs on cheap CPU machines — orders of magnitude cheaper than GPU fine-tuning."],
            size=13)
add_footer(s)

# ============================================================
# Slide 9 — NQ-320K results
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "6a · Document Retrieval — NQ-320K",
           "ERMs are fine-tuned; PEFA applied on top",
           page=9, total=TOTAL)

data = [
    ["Method",                  "R@10",   "R@100"],
    ["BM-25",                   "32.48",  "50.54"],
    ["DSI (base)",              "56.60",  "—"],
    ["SEAL (large)",            "81.24",  "90.93"],
    ["NCI (base)  — prev. SoTA","85.20",  "92.42"],
    ["MPNet_base",              "80.82",  "92.39"],
    ["  + PEFA-XS",             "86.67",  "94.53"],
    ["  + PEFA-XL",             "88.72",  "95.13"],
    ["GTR_base",                "79.74",  "90.91"],
    ["  + PEFA-XS",             "84.90",  "93.28"],
    ["  + PEFA-XL",             "88.71",  "94.36"],
    ["Avg. gain  PEFA-XS",      "+9.22",  "+5.28"],
    ["Avg. gain  PEFA-XL",      "+11.82", "+5.72"],
]
add_table(s, Inches(0.5), Inches(1.4), Inches(6.5), Inches(5.5),
          data, body_size=11, header_size=12,
          col_widths=[3.4, 1.55, 1.55])

# Right callouts
add_text(s, Inches(7.3), Inches(1.4), Inches(5.6), Inches(0.5),
         "Headline findings",
         size=18, bold=True, color=PRIMARY)
add_bullets(s, Inches(7.3), Inches(1.95), Inches(5.6), Inches(2.6),
            ["MPNet + PEFA-XL → 88.72 / 95.13  ⇒ new SoTA, beats NCI (85.20 / 92.42)",
             "GTR + PEFA-XL nearly ties (88.71)",
             "PEFA-XL > PEFA-XS on R@10; both >> baseline ERM"],
            size=13)

add_callout(s, Inches(7.3), Inches(4.7), Inches(5.6), Inches(2.0),
            "Why this matters",
            "Sequence-to-sequence indexers (NCI, DSI, SEAL) need expensive "
            "decoder training. PEFA reaches SoTA on dense bi-encoders "
            "with no extra training — only an HNSW build.")
add_footer(s)

# ============================================================
# Slide 10 — Trivia-QA + Ablation
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "6b · Trivia-QA  +  Hyper-parameter Ablation",
           "Robustness on pre-trained-only ERMs · effect of λ and k",
           page=10, total=TOTAL)

# Left: Trivia-QA table (pre-trained ERMs)
add_text(s, Inches(0.5), Inches(1.4), Inches(6.0), Inches(0.4),
         "Trivia-QA — ERMs are pre-trained only (no FT)",
         size=14, bold=True, color=PRIMARY)
data = [
    ["Method",                  "R@20",  "R@100"],
    ["NCI (base, FT)",          "94.45", "96.94"],
    ["Sent-BERT_distill",       "51.94", "68.50"],
    ["  + PEFA-XS",             "86.28", "93.33"],
    ["  + PEFA-XL",             "83.76", "91.83"],
    ["GTR_base",                "71.75", "82.05"],
    ["  + PEFA-XS",             "83.81", "91.02"],
    ["  + PEFA-XL",             "85.30", "92.38"],
    ["Avg. gain  PEFA-XS",      "+18.67","+13.61"],
    ["Avg. gain  PEFA-XL",      "+17.07","+12.80"],
]
add_table(s, Inches(0.5), Inches(1.85), Inches(6.0), Inches(4.4),
          data, body_size=11, header_size=12,
          col_widths=[3.0, 1.5, 1.5])

# Right: ablation insights
add_text(s, Inches(6.9), Inches(1.4), Inches(6.0), Inches(0.4),
         "Ablation on NQ-320K  (R@100)",
         size=14, bold=True, color=PRIMARY)
data2 = [
    ["ERM / variant", "λ=0.1", "λ=0.3", "λ=0.5", "λ=0.7", "λ=0.9"],
    ["GTR  PEFA-XS",     "92.11", "93.07", "93.31", "92.85", "91.74"],
    ["GTR  PEFA-XL k=16","94.36", "93.32", "92.81", "92.53", "91.93"],
    ["GTR  PEFA-XL k=32","94.32", "93.23", "92.82", "92.44", "91.79"],
    ["GTR  PEFA-XL k=64","93.93", "93.14", "92.76", "92.29", "91.62"],
    ["DPR  PEFA-XS",     "91.48", "92.22", "91.71", "89.87", "87.08"],
    ["DPR  PEFA-XL k=32","92.07", "90.50", "89.20", "88.62", "87.46"],
]
add_table(s, Inches(6.9), Inches(1.85), Inches(6.0), Inches(2.8),
          data2, body_size=10, header_size=11,
          col_widths=[2.4, 0.72, 0.72, 0.72, 0.72, 0.72])

add_text(s, Inches(6.9), Inches(4.8), Inches(6.0), Inches(0.4),
         "Take-aways",
         size=13, bold=True, color=ACCENT)
add_bullets(s, Inches(6.9), Inches(5.15), Inches(6.0), Inches(1.9),
            ["PEFA-XS optimum ≈ λ = 0.3 – 0.5 ; PEFA-XL ≈ λ = 0.1",
             "k = 32 typically saturates",
             "λ = 1 reduces to the original ERM (sanity check)"],
            size=12)
add_footer(s)

# ============================================================
# Slide 11 — Product search recall
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "7a · Product Search — Recall Gains",
           "Three proprietary datasets · numbers = absolute gain over MPNet baseline",
           page=11, total=TOTAL)

data = [
    ["Method", "5M R@100", "5M R@1000", "15M R@100", "15M R@1000", "30M R@100", "30M R@1000"],
    ["MPNet_base",        "0.00",  "0.00",  "0.00",  "0.00",  "0.00",  "0.00"],
    ["  + PEFA-XS",       "11.23", "13.14", "5.05",  "11.79", "9.67",  "17.47"],
    ["  + PEFA-XL",       "22.83", "12.31", "23.48", "21.56", "27.22", "18.96"],
    ["GTR_base",          "7.85",  "9.23",  "6.75",  "10.33", "8.35",  "9.83"],
    ["  + PEFA-XS",       "17.32", "19.55", "16.83", "25.00", "18.49", "24.38"],
    ["  + PEFA-XL",       "27.79", "19.23", "27.87", "28.75", "31.71", "24.28"],
    ["E5_base",           "9.93",  "9.75",  "9.98",  "12.98", "12.01", "12.61"],
    ["  + PEFA-XS",       "19.23", "19.18", "17.21", "27.78", "20.11", "26.08"],
    ["  + PEFA-XL",       "26.83", "17.75", "30.48", "31.07", "31.91", "25.49"],
    ["FT-ERM  (full FT)", "21.32", "20.87", "21.74", "30.04", "18.49", "24.11"],
    ["  + PEFA-XS",       "23.42", "22.17", "26.34", "34.84", "23.79", "29.61"],
    ["  + PEFA-XL",       "29.32", "22.87", "36.54", "37.24", "32.99", "30.01"],
]
add_table(s, Inches(0.4), Inches(1.45), Inches(12.6), Inches(5.0),
          data, body_size=10, header_size=10.5,
          col_widths=[2.8, 1.6, 1.6, 1.6, 1.6, 1.6, 1.6])

add_callout(s, Inches(0.4), Inches(6.55), Inches(12.6), Inches(0.6),
            "Highlight",
            "Even on top of a heavily fine-tuned FT-ERM, PEFA-XL adds +32.99 R@100 / +30.01 R@1000 "
            "on ProdSearch-30M — without touching ERM weights.")
add_footer(s)

# ============================================================
# Slide 12 — Deployment cost (Index/latency)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "7b · Product Search — Deployment Cost",
           "Index size (GiB) · Build time (hr) · Latency (ms/query, single thread)",
           page=12, total=TOTAL)

data = [
    ["Dataset",         "Method",       "Index disk (GiB)", "Build time (hr)", "Latency (ms/q)"],
    ["ProdSearch-5M",   "FT-ERM",       "13.1",  "0.3", "0.82"],
    ["ProdSearch-5M",   "+ PEFA-XS",    "13.1",  "0.2", "0.67"],
    ["ProdSearch-5M",   "+ PEFA-XL",    "32.2",  "0.7", "2.15"],
    ["ProdSearch-15M",  "FT-ERM",       "28.6",  "0.6", "0.91"],
    ["ProdSearch-15M",  "+ PEFA-XS",    "28.6",  "0.5", "0.94"],
    ["ProdSearch-15M",  "+ PEFA-XL",    "100.7", "1.9", "1.94"],
    ["ProdSearch-30M",  "FT-ERM",       "51.9",  "0.9", "0.77"],
    ["ProdSearch-30M",  "+ PEFA-XS",    "51.9",  "1.0", "0.71"],
    ["ProdSearch-30M",  "+ PEFA-XL",    "287.7", "4.7", "1.99"],
]
add_table(s, Inches(0.5), Inches(1.45), Inches(7.5), Inches(4.6),
          data, body_size=11, header_size=12,
          col_widths=[2.0, 1.7, 1.5, 1.4, 1.4])

# Right: insight column
add_text(s, Inches(8.3), Inches(1.45), Inches(4.7), Inches(0.4),
         "PEFA-XS",
         size=16, bold=True, color=PRIMARY)
add_bullets(s, Inches(8.3), Inches(1.85), Inches(4.7), Inches(1.7),
            ["Same index size as FT-ERM",
             "Latency unchanged (sometimes faster)",
             "Drop-in into production pipeline"],
            size=12)

add_text(s, Inches(8.3), Inches(3.7), Inches(4.7), Inches(0.4),
         "PEFA-XL",
         size=16, bold=True, color=ACCENT)
add_bullets(s, Inches(8.3), Inches(4.1), Inches(4.7), Inches(2.0),
            ["~3.6× index disk vs FT-ERM (30M)",
             "Latency ~2.4× (still < 2 ms/query)",
             "4.7 h CPU build vs hundreds of GPU-hours for FT",
             "Practitioner picks the trade-off"],
            size=12)
add_footer(s)

# ============================================================
# Slide 13 — Effect of supervised data  (figure3.png)
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "7c · Effect of Supervised Data",
           "How much of Y do we actually need? (Section 5.4)",
           page=13, total=TOTAL)

# Left: paper figure 3
add_rect(s, Inches(0.5), Inches(1.4), Inches(7.0), Inches(5.5),
         RGBColor(0xFF, 0xFF, 0xFF), line=PRIMARY)
add_text(s, Inches(0.5), Inches(1.45), Inches(7.0), Inches(0.35),
         "Figure 3 · ΔRecall@100 and index size vs. sampling ratio",
         size=12, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
s.shapes.add_picture(FIG3, Inches(0.7), Inches(1.85),
                     width=Inches(6.6), height=Inches(4.9))

# Right: narrative
add_text(s, Inches(7.8), Inches(1.4), Inches(5.2), Inches(0.4),
         "Findings",
         size=18, bold=True, color=PRIMARY)
add_bullets(s, Inches(7.8), Inches(1.9), Inches(5.2), Inches(2.8),
            ["PEFA-XS reaches FT-ERM-level R@100 with only ~10% of pairs sampled from Y",
             "Index size of PEFA-XS does NOT grow with more supervised data",
             "PEFA-XL achieves large gains even with just 5% of Y",
             "Sampling ratios tested: {0.05, 0.10, 0.25, 0.50, 0.75, 0.95}"],
            size=12)

add_callout(s, Inches(7.8), Inches(4.85), Inches(5.2), Inches(2.0),
            "Cost trade-off",
            "PEFA-XL index is ~1.4× larger than FT-ERM/PEFA-XS even at 5% sampling. "
            "Latency stays ~2× of FT-ERM across all dataset sizes. "
            "Practitioner picks Recall gain vs. infra cost.",
            color=ACCENT)
add_footer(s)

# ============================================================
# Slide 14 — Related work / Discussion
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "8 · Related Work & Discussion",
           "Where PEFA sits in the literature",
           page=14, total=TOTAL)

cards = [
    ("kNN-LM",
     "kNN component over vocabulary in language modeling. PEFA-XL is similar in spirit but operates on both passage AND query spaces; the latter is large and requires ANN, unlike a small vocab.",
     PRIMARY),
    ("PIFA (XMC)",
     "Positive Instance Feature Aggregation aggregates sparse tfidf features for clustering. PEFA-XS interpolates dense aggregated embeddings with the original passage embedding and indexes the result.",
     ACCENT),
    ("Parameter-efficient FT",
     "Adapter / LoRA / prefix-tuning still need gradients and underperform on retrieval (Tam et al.). PEFA needs no gradient → applies to black-box LLM encoders; orthogonal to PEFT.",
     GREEN),
    ("Ethical aspects",
     "Interpretability: kNN reveals which training queries drove a retrieval result.  Privacy: yearly-aggregated, anonymized search logs — no single-session traceability.",
     ORANGE),
]
x0 = Inches(0.5)
y0 = Inches(1.5)
cw = Inches(6.05)
ch = Inches(2.55)
gap = Inches(0.15)
for i, (title, body, col) in enumerate(cards):
    r, c = i // 2, i % 2
    x = x0 + c * (cw + gap)
    y = y0 + r * (ch + gap)
    add_rect(s, x, y, cw, ch, RGBColor(0xFF, 0xFF, 0xFF), line=col)
    add_rect(s, x, y, Inches(0.18), ch, col)
    add_text(s, x + Inches(0.3), y + Inches(0.1),
             cw - Inches(0.4), Inches(0.5),
             title, size=16, bold=True, color=col)
    add_text(s, x + Inches(0.3), y + Inches(0.55),
             cw - Inches(0.4), ch - Inches(0.6),
             body, size=12, color=DARK_TXT)
add_footer(s)

# ============================================================
# Slide 15 — Conclusion & My take
# ============================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Conclusion  &  My Take",
           "Summary, strengths, limitations",
           page=15, total=TOTAL)

# Left: paper conclusions
add_text(s, Inches(0.5), Inches(1.4), Inches(6.2), Inches(0.4),
         "Paper conclusions",
         size=16, bold=True, color=PRIMARY)
add_bullets(s, Inches(0.5), Inches(1.85), Inches(6.2), Inches(4.5),
            ["Parameter-free, gradient-free adapter framework for ERMs",
             "PEFA-XS: zero deployment overhead, modest gain",
             "PEFA-XL: bigger gain at ~2× latency and ~3.6× index",
             "Works on both pre-trained and fine-tuned ERMs",
             "NQ-320K: MPNet + PEFA-XL = 88.72 R@10 → new SoTA",
             "ProdSearch-30M: +5.3% / +14.5% R@100 over FT-ERM"],
            size=13)

# Right: my take
add_rect(s, Inches(7.0), Inches(1.4), Inches(5.9), Inches(5.5),
         SOFT_BG)
add_text(s, Inches(7.2), Inches(1.5), Inches(5.6), Inches(0.4),
         "My critical take",
         size=16, bold=True, color=ACCENT)
add_bullets(s, Inches(7.2), Inches(1.95), Inches(5.6), Inches(2.6),
            ["Strength · Reframes adaptation as ANN engineering — leverages mature CPU infra",
             "Strength · Inherently more interpretable than monolithic ERMs",
             "Limit · PEFA-XL needs label-side data Y at query time → not zero-shot",
             "Limit · Storage of Q + Y can dominate for very sparse Y",
             "Open · Combination with re-ranking; query-side augmentation; streaming Y updates"],
            size=12)
add_footer(s)

# ---------- Save ----------
prs.save(OUT_PATH)
print("Saved:", OUT_PATH)
print("Slides:", len(prs.slides))
